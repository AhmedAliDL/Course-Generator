import imghdr
import tempfile
import uuid
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont
import requests
from bing_image_urls import bing_image_urls
import os
import re


class GrapImage:

    # Function to sanitize the filename
    @staticmethod
    def sanitize_filename(filename):
        """
            Sanitizes a given filename by replacing invalid characters and ensuring it's safe to use in the filesystem.

            Parameters:
            - filename (str): The original filename to sanitize.

            Returns:
            - str: The sanitized filename.
            """
        # Replace invalid characters with underscores
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        # Replace any remaining non-alphanumeric characters (except . and _) with an underscore
        filename = re.sub(r'[^a-zA-Z0-9._]', '_', filename)
        # Remove any leading or trailing whitespace
        filename = filename.strip()
        # Ensure filename is not empty
        if not filename:
            filename = 'default_filename'
        return filename

    # Function to download and save image
    def download_and_save_image(self, image_content, title):
        """
           Downloads and saves an image from its content, assigning a filename based on the provided title.

           Parameters:
           - image_content (bytes): The binary content of the image to save.
           - title (str): The title to use for generating the filename.

           Returns:
           - str: The path where the image was saved, or None if saving failed.
           """
        try:
            if image_content is None:
                print("Image content is None. Skipping image.")
                return None

            sanitized_title = self.sanitize_filename(title)

            # Save the image content to a temporary file
            with tempfile.NamedTemporaryFile(delete=False) as temp:
                temp.write(image_content)
                temp_path = temp.name

            # Determine the file extension using imghdr
            file_extension = imghdr.what(temp_path)

            if file_extension is None or file_extension.lower() == 'webp':
                print("Unsupported image format (WEBP). Skipping image.")
                return None

            img_path = os.path.join("images", f"{sanitized_title}.{file_extension.lower()}")
            with open(img_path, "wb") as img_file:
                img_file.write(image_content)
            return img_path
        except Exception as e:
            print(f"Could not save image. Reason: {e}")
            return None

    # Function to scrape images from Bing
    @staticmethod
    def scrape_bing_images(data):
        """
            Scrapes Bing for images related to the topics provided in the data dictionary.

            Parameters:
            - data (dict): A dictionary where keys are lecture names and values are lists of topics.

            Returns:
            - dict: A dictionary mapping lecture names to another dictionary of topic-image URL pairs.
            """
        final_images = {}
        for lec, topics in data.items():
            images = {}
            for topic in topics:
                image_urls = bing_image_urls(topic, limit=3)
                if image_urls:
                    for url in image_urls:
                        if url.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.bmp')):
                            images[topic] = url
                            break  # Select the first valid image URL
            final_images[lec] = images
        return final_images


class Font:

    # Function to generate font preview image
    @staticmethod
    def generate_font_preview(font_name, font_size):
        """
           Generates a preview image of a specified font and size.

           Parameters:
           - font_name (str): The name of the font to preview.
           - font_size (int): The size of the font in points.

           Returns:
           - Image: A PIL Image object showing a sample text in the specified font and size,
           or None if the font is unavailable.
           """
        font_files = {
            "Arial": "arial.ttf",
            "Calibri": "calibri.ttf",
            "Times New Roman": "times.ttf",
            "Verdana": "verdana.ttf",
            "Courier New": "cour.ttf",
            "Georgia": "georgia.ttf",
            "Tahoma": "tahoma.ttf",
            "Comic Sans MS": "comic.ttf",
            "Impact": "impact.ttf",
            "Trebuchet MS": "trebuc.ttf"
        }

        if font_name in font_files:
            font_path = os.path.join(os.path.dirname(__file__), "fonts", font_files[font_name])
            image = Image.new("RGB", (300, 116), color="white")
            draw = ImageDraw.Draw(image)
            font = ImageFont.truetype(font_path, size=font_size)
            text = f"{font_name} 1 2 3"
            text_bbox = draw.textbbox((0, 0), text, font=font)
            text_width = text_bbox[2] - text_bbox[0]
            text_height = text_bbox[3] - text_bbox[1]
            draw.text(((300 - text_width) / 2, (100 - text_height) / 2), text, fill="black", font=font)
            return image
        else:
            return None


class CreatePresentation:
    def __init__(self):
        self.grapImage = GrapImage()

    # Function to add a slide with image
    def add_slide_with_image(self, presentation, title_text, content_text, image_url, title_font_size,
                             content_font_size,
                             title_color, content_color, font_name, background_image_path=None,
                             is_lecture=False):
        """
           Adds a slide to the presentation with specified title, content, and an optional image.

           Parameters:
           - presentation (Presentation): The Presentation object to add slides to.
           - title_text (str): The text for the slide's title.
           - content_text (str): The main content text for the slide.
           - image_url (str): The URL of the image to include in the slide.
           - title_font_size (int): The font size for the title.
           - content_font_size (int): The font size for the content.
           - title_color (RGBColor): The color of the title text.
           - content_color (RGBColor): The color of the content text.
           - font_name (str): The font name for both title and content.
           - background_image_path (str, optional): The path to a background image for the slide. Defaults to None.
           - is_lecture (bool): Flag indicating if this is a lecture slide. Defaults to False.
           """
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        if background_image_path:
            slide.shapes.add_picture(background_image_path, 0, 0, width=presentation.slide_width,
                                     height=presentation.slide_height)

        # Add title and content to the slide
        title_left = Inches(0.2)
        title_width = presentation.slide_width - Inches(2)
        title_height = Inches(1.1)
        if is_lecture:
            title_font_size = 60
            title_left = Inches(3)
            title_top = (presentation.slide_height - title_height) / 2
        else:
            title_top = Inches(0.5)

        title_text_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_text_box.text_frame
        title_frame.word_wrap = True
        title_paragraph = title_frame.add_paragraph()
        title_paragraph.text = title_text
        title_paragraph.font.size = Pt(title_font_size)
        title_paragraph.font.color.rgb = title_color
        title_paragraph.font.name = font_name

        # Add content to the slide

        # Add image if available
        if image_url:
            try:
                img_response = requests.get(image_url)
                img_response.raise_for_status()
                img_path = self.grapImage.download_and_save_image(img_response.content, title_text)
                if img_path:
                    slide.shapes.add_picture(img_path, left=Inches(6), top=Inches(2.25), width=Inches(4),
                                             height=Inches(4))
                    content_left = Inches(0.6)
                    content_top = Inches(1.5)
                    content_width = presentation.slide_width - Inches(5)
                    content_height = presentation.slide_height - Inches(2.5)
                    content_text_box = slide.shapes.add_textbox(content_left, content_top, content_width,
                                                                content_height)
                    content_frame = content_text_box.text_frame
                    content_frame.word_wrap = True
                    content_paragraph = content_frame.add_paragraph()
                    content_paragraph.text = content_text
                    content_paragraph.font.size = Pt(content_font_size)
                    content_paragraph.font.color.rgb = content_color
                    content_paragraph.font.name = font_name
                else:

                    content_left = Inches(0.6)
                    content_top = Inches(1.5)
                    content_width = presentation.slide_width - Inches(1)
                    content_height = presentation.slide_height - Inches(2.5)
                    content_text_box = slide.shapes.add_textbox(content_left, content_top, content_width,
                                                                content_height)
                    content_frame = content_text_box.text_frame
                    content_frame.word_wrap = True
                    content_paragraph = content_frame.add_paragraph()
                    content_paragraph.text = content_text
                    content_paragraph.font.size = Pt(content_font_size)
                    content_paragraph.font.color.rgb = content_color
                    content_paragraph.font.name = font_name
            except requests.exceptions.RequestException as e:
                print(f"Could not download {image_url}. Reason: {e}")
                content_left = Inches(0.6)
                content_top = Inches(1.5)
                content_width = presentation.slide_width - Inches(1)
                content_height = presentation.slide_height - Inches(2.5)
                content_text_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                content_frame = content_text_box.text_frame
                content_frame.word_wrap = True
                content_paragraph = content_frame.add_paragraph()
                content_paragraph.text = content_text
                content_paragraph.font.size = Pt(content_font_size)
                content_paragraph.font.color.rgb = content_color
                content_paragraph.font.name = font_name
        else:
            content_left = Inches(0.6)
            content_top = Inches(1.5)
            content_width = presentation.slide_width - Inches(1)
            content_height = presentation.slide_height - Inches(2.5)
            content_text_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_text_box.text_frame
            content_frame.word_wrap = True
            content_paragraph = content_frame.add_paragraph()
            content_paragraph.text = content_text
            content_paragraph.font.size = Pt(content_font_size)
            content_paragraph.font.color.rgb = content_color
            content_paragraph.font.name = font_name

    def create_presentation(self, optimizedPres, title_font_size, content_font_size,
                            title_color, content_color, font_name, background_image_path, name):
        """
            Creates a PowerPoint presentation based on provided lecture data and user-specified styles.

            Parameters:
            - lectures_data (dict): A dictionary where keys are lecture titles and values are dictionaries of topics with content.
            - name (str): The name to use for saving the presentation file.

            Returns:
            - None: This function saves the presentation file to the local filesystem and provides an option to download it.
            """

        try:
            with st.spinner("create powerpoint"):
                if "file_path" not in st.session_state:
                    st.session_state.file_path = None

                if "file_name" not in st.session_state:
                    st.session_state.file_name = None

                prs = Presentation()

                image_urls = GrapImage.scrape_bing_images(optimizedPres)

                for lecture_title, topics_data in optimizedPres.items():
                    self.add_slide_with_image(prs, lecture_title, "", None, title_font_size, content_font_size,
                                              title_color, content_color, font_name,
                                              background_image_path=background_image_path,
                                              is_lecture=True)

                    for topic_title, content_lines in topics_data.items():
                        # Ensure content_lines is always a list
                        if isinstance(content_lines, str):
                            content_lines = [content_lines]

                        # Split content into chunks of 3 lines or fewer
                        for i in range(0, len(content_lines), 3):
                            chunk = content_lines[i:i + 3]
                            content_text = "\n\n".join(chunk)
                            self.add_slide_with_image(prs, topic_title, content_text,
                                                      image_urls.get(lecture_title, {}).get(topic_title),
                                                      title_font_size, content_font_size, title_color, content_color,
                                                      font_name,
                                                      background_image_path=background_image_path,
                                                      )

                st.session_state.file_name = name + '_' + str(uuid.uuid4())
                st.session_state.file_path = fr"Pres\powerpoint_{st.session_state.file_name}.pptx"
                prs.save(st.session_state.file_path)
                st.success("PowerPoint created successfully!")

        except Exception as e:
            st.error(f"An error occurred: {e}")
